from pptx import Presentation
import os

def openPPT(filename):
	os.system("\"C:\\Program Files\\Microsoft Office\\Office14\\POWERPNT.EXE\" "+ filename)




print("holi")
prs = Presentation('lala.pptx')
#f = open('test.pptx')

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "FALAMOS"
subtitle.text = "python-pptx was here!"

prs.save("lala.pptx")
openPPT("lala.pptx")
